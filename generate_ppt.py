"""
Generate ICSSCS 2026 Conference Presentation
Handwritten Digit Recognition Using Convolutional Neural Networks: A Deep Learning Approach

Uses the official ICSSCS 2026 PPT Format template as the base.
Sections: Introduction, Related Work, Methodology, Result Analysis, Conclusion, References
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import copy
import os

# --- Configuration ---
TEMPLATE_PATH = r"c:\HandWrittenDigit\ICSSCS 2026 PPT Format.pptx"
OUTPUT_PATH = r"c:\HandWrittenDigit\ICSSCS_2026_Presentation_v2.pptx"

# Colors from the template
COLOR_TITLE = RGBColor(0xA5, 0x30, 0x10)
COLOR_GREEN = RGBColor(0x72, 0x86, 0x53)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
COLOR_ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
COLOR_LIGHT_GREEN = RGBColor(0x4C, 0x8C, 0x4A)
COLOR_DARK_RED = RGBColor(0x8B, 0x1A, 0x1A)
COLOR_TABLE_HEADER = RGBColor(0x2E, 0x4A, 0x6E)
COLOR_TABLE_ROW_ALT = RGBColor(0xE8, 0xEF, 0xF5)
COLOR_HIGHLIGHT = RGBColor(0xC0, 0x39, 0x2B)

FONT_TITLE = "Times New Roman"
FONT_BODY = "Century Gothic"


# --- Helper Functions ---

def add_textbox(slide, left, top, width, height, text="", font_name=FONT_BODY,
                font_size=14, font_color=COLOR_BLACK, bold=False, italic=False,
                alignment=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    if text:
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.italic = italic
    return tf


def add_paragraph(text_frame, text="", font_name=FONT_BODY, font_size=14,
                  font_color=COLOR_BLACK, bold=False, italic=False,
                  alignment=PP_ALIGN.LEFT, space_before=0, space_after=0):
    p = text_frame.add_paragraph()
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if text:
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.italic = italic
    return p


def add_slide_title(slide, title_text, left=None, top=None, width=None, height=None):
    if left is None: left = 2800000
    if top is None: top = 100000
    if width is None: width = 3500000
    if height is None: height = 650000
    tf = add_textbox(slide, left, top, width, height,
                     text=title_text, font_name=FONT_TITLE,
                     font_size=32, font_color=COLOR_DARK_BLUE,
                     bold=True, alignment=PP_ALIGN.CENTER)
    return tf


def add_bullet_points(slide, points, left, top, width, height,
                      font_size=13, font_color=COLOR_BLACK, bullet_char="\u2022"):
    tf = add_textbox(slide, left, top, width, height, font_name=FONT_BODY, font_size=font_size)
    first = True
    for point in points:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        p.space_after = Pt(3)
        run_bullet = p.add_run()
        run_bullet.text = f"{bullet_char} "
        run_bullet.font.name = FONT_BODY
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = COLOR_HIGHLIGHT
        run_bullet.font.bold = True
        run_text = p.add_run()
        run_text.text = point
        run_text.font.name = FONT_BODY
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = font_color
    return tf


def add_table(slide, rows_data, left, top, width, height, col_widths=None):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 1
    table_shape = slide.shapes.add_table(n_rows, n_cols, Emu(left), Emu(top), Emu(width), Emu(height))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(w)
    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(cell_text)
            if row_idx == 0:
                run.font.name = FONT_BODY
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = COLOR_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TABLE_HEADER
            else:
                run.font.name = FONT_BODY
                run.font.size = Pt(9)
                run.font.color.rgb = COLOR_BLACK
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TABLE_ROW_ALT if row_idx % 2 == 0 else COLOR_WHITE
            cell.margin_left = Emu(45720)
            cell.margin_right = Emu(45720)
            cell.margin_top = Emu(18288)
            cell.margin_bottom = Emu(18288)
    return table


def add_decorative_line(slide, left, top, width, color=COLOR_TITLE, weight=2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Pt(weight))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_content_slide(prs, template_slide):
    layout = template_slide.slide_layout
    slide = prs.slides.add_slide(layout)
    return slide


# --- Main Generation ---

def generate_presentation():
    prs = Presentation(TEMPLATE_PATH)
    content_slide_ref = prs.slides[2]

    # ============================================================
    # SLIDE 1: TITLE SLIDE (modify existing)
    # ============================================================
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame and "Paper ID" in shape.text:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            tf.paragraphs[0].runs[0].text = "Paper ID: HDR-CNN-2026"
            tf.paragraphs[0].runs[0].font.name = FONT_BODY
            tf.paragraphs[0].runs[0].font.size = Pt(15.6)
            tf.paragraphs[0].runs[0].font.color.rgb = COLOR_GREEN
            p_title = tf.add_paragraph()
            p_title.alignment = PP_ALIGN.CENTER
            run_t = p_title.add_run()
            run_t.text = "Handwritten Digit Recognition Using Convolutional Neural Networks:"
            run_t.font.name = FONT_BODY
            run_t.font.size = Pt(15.6)
            run_t.font.color.rgb = COLOR_GREEN
            p_title2 = tf.add_paragraph()
            p_title2.alignment = PP_ALIGN.CENTER
            run_t2 = p_title2.add_run()
            run_t2.text = "A Deep Learning Approach"
            run_t2.font.name = FONT_BODY
            run_t2.font.size = Pt(15.6)
            run_t2.font.color.rgb = COLOR_GREEN

        if shape.has_text_frame and "Presented by" in shape.text:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            tf.paragraphs[0].runs[0].text = "Presented by"
            tf.paragraphs[0].runs[0].font.name = FONT_BODY
            tf.paragraphs[0].runs[0].font.size = Pt(14)
            tf.paragraphs[0].runs[0].font.bold = True
            tf.paragraphs[0].runs[0].font.color.rgb = COLOR_BLACK
            tf.paragraphs[1].runs[0].text = "Abhay Singh"
            tf.paragraphs[1].runs[0].font.name = FONT_BODY
            tf.paragraphs[1].runs[0].font.size = Pt(14)
            tf.paragraphs[1].runs[0].font.bold = True
            tf.paragraphs[1].runs[0].font.color.rgb = COLOR_BLACK
            tf.paragraphs[2].runs[0].text = "Department of Computer Science, COER University, Roorkee, India"
            tf.paragraphs[2].runs[0].font.name = FONT_BODY
            tf.paragraphs[2].runs[0].font.size = Pt(14)
            tf.paragraphs[2].runs[0].font.bold = True
            tf.paragraphs[2].runs[0].font.color.rgb = COLOR_BLACK
            tf.paragraphs[4].runs[0].text = "Presentation Date: April 2026"
            tf.paragraphs[4].runs[0].font.name = FONT_BODY
            tf.paragraphs[4].runs[0].font.size = Pt(14)
            tf.paragraphs[4].runs[0].font.bold = True
            tf.paragraphs[4].runs[0].font.color.rgb = COLOR_BLACK

    # ============================================================
    # SLIDE 2: TABLE OF CONTENTS (modify existing)
    # ============================================================
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_text_frame and "Introduction" in shape.text:
            tf = shape.text_frame
            toc_items = [
                "Introduction",
                "Related Work",
                "Methodology",
                "Result Analysis",
                "Conclusion",
                "References"
            ]
            for i, para in enumerate(tf.paragraphs):
                if i < len(toc_items):
                    for run in para.runs:
                        run.text = ""
                    para.runs[0].text = f"{i+1}.  {toc_items[i]}"
                    para.runs[0].font.name = FONT_BODY
                    para.runs[0].font.size = Pt(15.6)
                    para.runs[0].font.bold = True
                else:
                    # Clear extra paragraphs
                    for run in para.runs:
                        run.text = ""

    # ============================================================
    # SLIDE 3: INTRODUCTION (modify existing template slide 3)
    # ============================================================
    slide3 = prs.slides[2]
    # The title "Introduction" is already there from the template

    add_textbox(slide3, 300000, 850000, 4300000, 350000,
                text="Background & Problem",
                font_name=FONT_BODY, font_size=16, font_color=COLOR_DARK_RED, bold=True)

    add_bullet_points(slide3, [
        "Handwritten digit recognition (HDR) is a key problem in pattern recognition & computer vision",
        "Practical uses: postal mail sorting, bank cheque processing, document scanning, mobile input",
        "Traditional methods (SVM, k-NN) rely on manual feature extraction and miss spatial patterns",
        "CNNs learn features automatically from raw pixels at multiple levels of abstraction",
    ], left=300000, top=1150000, width=4300000, height=2200000, font_size=12)

    add_textbox(slide3, 4800000, 850000, 4100000, 350000,
                text="Objectives",
                font_name=FONT_BODY, font_size=16, font_color=COLOR_DARK_RED, bold=True)

    add_bullet_points(slide3, [
        "Design and build a CNN model for digit classification on MNIST dataset",
        "Measure performance: accuracy, loss, and generalization capability",
        "Study the role of each building block (convolution, pooling, dropout, batch norm)",
        "Deploy trained model in a real-time Streamlit web application",
        "Suggest improvements and discuss real-world applications",
    ], left=4800000, top=1150000, width=4100000, height=2600000, font_size=12)

    # Abstract summary at bottom
    add_textbox(slide3, 300000, 3800000, 8500000, 1000000,
                text="Abstract: This paper presents a CNN-based approach using the MNIST benchmark dataset. "
                     "The model uses two convolutional layers with ReLU, max-pooling, batch normalisation, and dropout. "
                     "Built with TensorFlow/Keras, trained on 60,000 images, achieving 98.91% test accuracy. "
                     "Includes real-time deployment through a Streamlit web application.",
                font_name=FONT_BODY, font_size=11, font_color=COLOR_DARK_BLUE,
                bold=False, italic=True, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 4: RELATED WORK (Slide 1 - Classical vs Modern)
    # ============================================================
    slide4 = add_content_slide(prs, content_slide_ref)
    add_slide_title(slide4, "Related Work")
    add_decorative_line(slide4, 3200000, 700000, 2700000, COLOR_TITLE)

    add_textbox(slide4, 300000, 850000, 4200000, 350000,
                text="Classical Approaches",
                font_name=FONT_BODY, font_size=16, font_color=COLOR_DARK_RED, bold=True)

    add_bullet_points(slide4, [
        "SVM (Support Vector Machines): ~1.4% error on MNIST [Vapnik, 1995]",
        "k-Nearest Neighbours (k-NN): 3-5% error depending on features used",
        "Random Forests: 96-97% accuracy on MNIST",
        "All required handcrafted features (HOG, Gabor filters)",
        "Could not capture spatial relationships in image data"
    ], left=300000, top=1150000, width=4200000, height=2000000, font_size=12)

    add_textbox(slide4, 4700000, 850000, 4200000, 350000,
                text="Deep Learning Approaches",
                font_name=FONT_BODY, font_size=16, font_color=COLOR_DARK_RED, bold=True)

    add_bullet_points(slide4, [
        "LeNet-5 (LeCun, 1998): 0.80% error - first successful CNN for digits",
        "Deep CNN (Cirisan, 2012): 0.35% error - deeper networks",
        "DropConnect (2013): 0.21% error - better regularisation",
        "Batch-Normalised CNN (2015): 0.29% error - stable training",
        "Ensemble of CNNs (2016): 0.17% error - combining models",
        "Capsule Network (Hinton, 2017): 0.25% - spatial relationships"
    ], left=4700000, top=1150000, width=4200000, height=2200000, font_size=12)

    # Comparison table
    add_table(slide4, [
        ["Method", "Error Rate (%)", "Year"],
        ["LeNet-5", "0.80", "1998"],
        ["SVM (RBF kernel)", "1.40", "1998"],
        ["Deep CNN", "0.35", "2012"],
        ["DropConnect", "0.21", "2013"],
        ["Batch-Norm CNN", "0.29", "2015"],
        ["Ensemble CNNs", "0.17", "2016"],
        ["Capsule Network", "0.25", "2017"],
    ], left=1500000, top=3400000, width=6100000, height=1500000)

    # Research gap
    add_textbox(slide4, 300000, 4850000, 8500000, 250000,
                text="Research Gap: Most works focus only on accuracy without practical deployment. "
                     "This paper provides reproducible implementation + real-time web application.",
                font_name=FONT_BODY, font_size=11, font_color=COLOR_HIGHLIGHT,
                bold=True, italic=True, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 5: METHODOLOGY (Slide 1 - Dataset & Preprocessing)
    # ============================================================
    slide5 = add_content_slide(prs, content_slide_ref)
    add_slide_title(slide5, "Methodology", width=4000000)
    add_decorative_line(slide5, 3000000, 700000, 3100000, COLOR_TITLE)

    add_textbox(slide5, 300000, 850000, 4200000, 350000,
                text="MNIST Dataset",
                font_name=FONT_BODY, font_size=16, font_color=COLOR_DARK_RED, bold=True)

    add_bullet_points(slide5, [
        "70,000 grayscale images (60K train + 10K test)",
        "28 x 28 pixels, white digit on black background",
        "10 classes (digits 0-9), roughly balanced",
        "Most widely used benchmark for digit recognition",
    ], left=300000, top=1150000, width=4200000, height=1200000, font_size=12)

    # Dataset table
    add_table(slide5, [
        ["Property", "Value"],
        ["Total Images", "70,000"],
        ["Training Set", "60,000"],
        ["Test Set", "10,000"],
        ["Dimensions", "28 x 28 pixels"],
        ["Colour Space", "Grayscale (1 channel)"],
        ["Classes", "10 (digits 0-9)"],
    ], left=4800000, top=850000, width=3800000, height=1600000)

    # Preprocessing
    add_textbox(slide5, 300000, 2600000, 8500000, 350000,
                text="Data Preprocessing Pipeline",
                font_name=FONT_BODY, font_size=16, font_color=COLOR_DARK_BLUE,
                bold=True, alignment=PP_ALIGN.CENTER)

    # Three steps side by side
    add_textbox(slide5, 300000, 2950000, 2700000, 300000,
                text="1. Reshaping",
                font_name=FONT_BODY, font_size=14, font_color=COLOR_DARK_RED, bold=True)
    add_bullet_points(slide5, [
        "(N, 28, 28) -> (N, 28, 28, 1)",
        "Add channel dimension for CNN"
    ], left=300000, top=3200000, width=2700000, height=600000, font_size=11)

    add_textbox(slide5, 3200000, 2950000, 2800000, 300000,
                text="2. Normalisation",
                font_name=FONT_BODY, font_size=14, font_color=COLOR_DARK_RED, bold=True)
    add_bullet_points(slide5, [
        "Pixel values / 255",
        "[0, 255] -> [0.0, 1.0]"
    ], left=3200000, top=3200000, width=2800000, height=600000, font_size=11)

    add_textbox(slide5, 6200000, 2950000, 2700000, 300000,
                text="3. One-Hot Encoding",
                font_name=FONT_BODY, font_size=14, font_color=COLOR_DARK_RED, bold=True)
    add_bullet_points(slide5, [
        "Label 3 -> [0,0,0,1,0,0,0,0,0,0]",
        "Vector of length 10 for loss fn"
    ], left=6200000, top=3200000, width=2700000, height=600000, font_size=11)

    # Pipeline flow
    add_textbox(slide5, 500000, 4000000, 8100000, 400000,
                text="Raw Images --> Reshape (28x28x1) --> Normalize (/255) --> One-Hot Encode --> CNN Input",
                font_name=FONT_BODY, font_size=13, font_color=COLOR_DARK_BLUE,
                bold=True, alignment=PP_ALIGN.CENTER)

    # Environment
    add_table(slide5, [
        ["Component", "Specification"],
        ["Language", "Python 3.13"],
        ["Framework", "TensorFlow 2.21 / Keras"],
        ["Hardware", "CPU-based training"],
        ["Training Time", "~3 min (9 epochs)"],
        ["Deployment", "Streamlit Web App"],
    ], left=2200000, top=4400000, width=4700000, height=700000)

    # ============================================================
    # SLIDE 6: METHODOLOGY (Slide 2 - CNN Architecture)
    # ============================================================
    slide6 = add_content_slide(prs, content_slide_ref)
    add_slide_title(slide6, "Methodology: CNN Architecture", width=5000000, left=2000000)
    add_decorative_line(slide6, 2400000, 700000, 4300000, COLOR_TITLE)

    # Architecture text-based diagram
    arch_tf = add_textbox(slide6, 300000, 850000, 4000000, 3500000,
                          text="Proposed CNN Architecture",
                          font_name=FONT_BODY, font_size=14, font_color=COLOR_DARK_RED, bold=True)

    arch_lines = [
        ("Input Layer:", "28 x 28 x 1"),
        ("", ""),
        ("-- Feature Extraction --", ""),
        ("Conv2D(32, 3x3, ReLU)", "-> 26 x 26 x 32"),
        ("Conv2D(64, 3x3, ReLU)", "-> 24 x 24 x 64"),
        ("MaxPooling2D(2x2)", "-> 12 x 12 x 64"),
        ("Dropout(0.25)", "-> 12 x 12 x 64"),
        ("", ""),
        ("-- Classification Head --", ""),
        ("Flatten()", "-> 9,216"),
        ("Dense(128, ReLU)", "-> 128"),
        ("BatchNormalization()", "-> 128"),
        ("Dropout(0.50)", "-> 128"),
        ("Dense(10, Softmax)", "-> 10 (output)"),
    ]

    for layer_name, output_shape in arch_lines:
        p = arch_tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run1 = p.add_run()
        run1.text = f"  {layer_name}"
        run1.font.name = "Consolas"
        run1.font.size = Pt(10)
        run1.font.bold = "--" in layer_name or layer_name == "Input Layer:"
        run1.font.color.rgb = COLOR_DARK_BLUE if "--" in layer_name else COLOR_BLACK
        if output_shape:
            run2 = p.add_run()
            run2.text = f"  {output_shape}"
            run2.font.name = "Consolas"
            run2.font.size = Pt(10)
            run2.font.color.rgb = COLOR_ACCENT_BLUE

    p_params = arch_tf.add_paragraph()
    p_params.space_before = Pt(6)
    run_p = p_params.add_run()
    run_p.text = "  Total: ~1,199,882 params | Trainable: ~1,199,370"
    run_p.font.name = FONT_BODY
    run_p.font.size = Pt(10)
    run_p.font.color.rgb = COLOR_HIGHLIGHT
    run_p.font.bold = True

    # Layer explanations
    add_textbox(slide6, 4500000, 850000, 4300000, 350000,
                text="Layer-by-Layer Explanation",
                font_name=FONT_BODY, font_size=14, font_color=COLOR_DARK_RED, bold=True)

    add_bullet_points(slide6, [
        "Conv2D: Slides 3x3 filters to detect patterns (edges -> shapes -> digits)",
        "ReLU: f(x) = max(0, x) - avoids vanishing gradient, fast computation",
        "MaxPooling: Takes max from 2x2 blocks, reduces size by 75%",
        "Dropout (25% & 50%): Randomly deactivates neurons to prevent overfitting",
        "BatchNorm: Normalizes layer inputs, stabilizes and speeds up training",
        "Softmax: Outputs probability distribution over 10 classes (sum = 1.0)"
    ], left=4500000, top=1200000, width=4300000, height=2200000, font_size=11)

    # Training config table
    add_textbox(slide6, 4500000, 3500000, 4300000, 300000,
                text="Training Hyperparameters",
                font_name=FONT_BODY, font_size=13, font_color=COLOR_DARK_BLUE, bold=True)

    add_table(slide6, [
        ["Parameter", "Value"],
        ["Optimiser", "Adam (lr=0.001)"],
        ["Loss Function", "Categorical Cross-Entropy"],
        ["Batch Size", "128"],
        ["Max Epochs", "15"],
        ["Validation Split", "10%"],
        ["Early Stopping", "patience=3"],
    ], left=4500000, top=3750000, width=4200000, height=1200000)

    # ============================================================
    # SLIDE 7: RESULT ANALYSIS (Slide 1 - Training & Test)
    # ============================================================
    slide7 = add_content_slide(prs, content_slide_ref)
    add_slide_title(slide7, "Result Analysis")
    add_decorative_line(slide7, 3000000, 700000, 3100000, COLOR_TITLE)

    # Training performance
    add_textbox(slide7, 300000, 850000, 4200000, 350000,
                text="Training Performance",
                font_name=FONT_BODY, font_size=16, font_color=COLOR_DARK_RED, bold=True)

    add_bullet_points(slide7, [
        "Fast initial learning: 95.4% accuracy after Epoch 1",
        "Steady improvement: 95.4% -> 99.3% over 9 epochs",
        "No overfitting: Train-validation gap < 0.5%",
        "Validation plateau at ~99.1-99.2% by Epoch 5",
        "Early stopping activated at Epoch 9 (patience=3)"
    ], left=300000, top=1200000, width=4200000, height=1800000, font_size=12)

    # Test results
    add_textbox(slide7, 4800000, 850000, 4000000, 350000,
                text="Test Set Evaluation",
                font_name=FONT_BODY, font_size=16, font_color=COLOR_DARK_RED, bold=True)

    add_table(slide7, [
        ["Metric", "Value"],
        ["Test Accuracy", "98.91%"],
        ["Test Loss", "0.0304"],
        ["Error Rate", "1.09%"],
        ["Misclassified", "~109 / 10,000"],
    ], left=4800000, top=1200000, width=3800000, height=1200000)

    # Performance comparison
    add_textbox(slide7, 300000, 3200000, 8500000, 350000,
                text="Performance Comparison with Baseline Methods",
                font_name=FONT_BODY, font_size=14, font_color=COLOR_DARK_BLUE,
                bold=True, alignment=PP_ALIGN.CENTER)

    add_table(slide7, [
        ["Method", "Test Accuracy (%)", "Parameters"],
        ["k-NN (baseline)", "96.9", "--"],
        ["SVM (RBF)", "98.6", "--"],
        ["MLP (784-256-128-10)", "97.8", "~236K"],
        ["Proposed CNN", "98.91", "~1.2M"],
        ["LeNet-5", "99.2", "~60K"],
        ["Deep CNN (Cirisan)", "99.65", "~10M"],
    ], left=1200000, top=3500000, width=6700000, height=1400000)

    # ============================================================
    # SLIDE 8: RESULT ANALYSIS (Slide 2 - Per-Class & Confusion)
    # ============================================================
    slide8 = add_content_slide(prs, content_slide_ref)
    add_slide_title(slide8, "Result Analysis: Per-Class Report", width=5000000, left=2000000)
    add_decorative_line(slide8, 2400000, 700000, 4300000, COLOR_TITLE)

    # Classification report table
    add_table(slide8, [
        ["Digit", "Precision", "Recall", "F1-Score", "Support"],
        ["0", "0.9888", "0.9929", "0.9908", "980"],
        ["1", "0.9938", "0.9947", "0.9943", "1,135"],
        ["2", "0.9799", "0.9942", "0.9870", "1,032"],
        ["3", "0.9891", "0.9911", "0.9901", "1,010"],
        ["4", "0.9928", "0.9898", "0.9913", "982"],
        ["5", "0.9757", "0.9910", "0.9833", "892"],
        ["6", "0.9947", "0.9802", "0.9874", "958"],
        ["7", "0.9903", "0.9893", "0.9898", "1,028"],
        ["8", "0.9938", "0.9877", "0.9907", "974"],
        ["9", "0.9910", "0.9792", "0.9850", "1,009"],
        ["Avg", "0.9891", "0.9891", "0.9891", "10,000"],
    ], left=300000, top=800000, width=5000000, height=3000000)

    # Key observations
    add_textbox(slide8, 5500000, 800000, 3400000, 350000,
                text="Key Observations",
                font_name=FONT_BODY, font_size=16, font_color=COLOR_DARK_RED, bold=True)

    add_bullet_points(slide8, [
        "Best recognized: Digit 1 (99.47% recall) - simple, straight stroke",
        "Most confused: Digit 9 (97.92% recall) - confused with digit 4",
        "Lowest precision: Digit 5 (97.57%) - similar curves to 3 & 6"
    ], left=5500000, top=1150000, width=3400000, height=1100000, font_size=12)

    add_textbox(slide8, 5500000, 2400000, 3400000, 350000,
                text="Common Misclassifications",
                font_name=FONT_BODY, font_size=14, font_color=COLOR_DARK_RED, bold=True)

    add_bullet_points(slide8, [
        "6 -> 5: 9 cases (similar curves)",
        "9 -> 4: 6 cases (similar top loops)",
        "9 -> 5: 6 cases (open top style)",
        "0 -> 2: 6 cases (slanted ovals)",
        "Most from genuinely ambiguous handwriting"
    ], left=5500000, top=2700000, width=3400000, height=1600000, font_size=11)

    # Computational efficiency
    add_textbox(slide8, 5500000, 4200000, 3400000, 300000,
                text="Computational Efficiency",
                font_name=FONT_BODY, font_size=13, font_color=COLOR_DARK_BLUE, bold=True)

    add_table(slide8, [
        ["Metric", "Value"],
        ["Training Time", "~3 min (CPU)"],
        ["Inference", "< 5 ms/image"],
        ["Model Size", "~14 MB"],
    ], left=5500000, top=4450000, width=3200000, height=550000)

    # ============================================================
    # SLIDE 9: CONCLUSION
    # ============================================================
    slide9 = add_content_slide(prs, content_slide_ref)
    add_slide_title(slide9, "Conclusion")
    add_decorative_line(slide9, 3300000, 700000, 2500000, COLOR_TITLE)

    # Conclusion summary
    add_textbox(slide9, 300000, 850000, 4200000, 350000,
                text="Summary of Findings",
                font_name=FONT_BODY, font_size=16, font_color=COLOR_DARK_RED, bold=True)

    add_bullet_points(slide9, [
        "CNN-based approach for handwritten digit recognition on MNIST",
        "Two conv layers + max-pooling + batch norm + dropout",
        "Achieved 98.91% test accuracy with a simple architecture",
        "Outperforms all traditional methods (k-NN, SVM, MLP)",
        "Complete end-to-end pipeline: training to web deployment"
    ], left=300000, top=1200000, width=4200000, height=1800000, font_size=13)

    # Key contributions
    add_textbox(slide9, 4800000, 850000, 4100000, 350000,
                text="Key Contributions",
                font_name=FONT_BODY, font_size=16, font_color=COLOR_DARK_RED, bold=True)

    add_bullet_points(slide9, [
        "Clearly documented, reproducible CNN implementation",
        "Detailed layer-by-layer explanation of architecture",
        "Thorough analysis of training behavior & regularization",
        "Working Streamlit web app for real-time recognition",
        "Honest discussion of limitations & improvements"
    ], left=4800000, top=1200000, width=4100000, height=1800000, font_size=13)

    # Future Work
    add_textbox(slide9, 300000, 3300000, 8500000, 350000,
                text="Future Work",
                font_name=FONT_BODY, font_size=16, font_color=COLOR_DARK_BLUE,
                bold=True, alignment=PP_ALIGN.CENTER)

    add_bullet_points(slide9, [
        "Multi-digit number recognition (e.g., \"142\") with digit segmentation",
        "Data augmentation: rotation, shifting, zoom to handle more writing variations",
        "Train on harder datasets: EMNIST (letters), SVHN (real-world street numbers)",
        "Deeper architecture with residual connections (ResNet) or Capsule Networks",
        "Model compression and knowledge distillation for mobile deployment"
    ], left=300000, top=3600000, width=8500000, height=1500000, font_size=12)

    # ============================================================
    # SLIDE 10: REFERENCES
    # ============================================================
    slide10 = add_content_slide(prs, content_slide_ref)
    add_slide_title(slide10, "References")
    add_decorative_line(slide10, 3300000, 700000, 2500000, COLOR_TITLE)

    refs_left = [
        "[1] R. Plamondon & S.N. Srihari, \"Online and off-line handwriting recognition,\" IEEE TPAMI, 2000.",
        "[2] Y. LeCun, Y. Bengio & G. Hinton, \"Deep learning,\" Nature, vol. 521, 2015.",
        "[3] I. Goodfellow et al., Deep Learning, MIT Press, 2016.",
        "[8] V. Vapnik, The Nature of Statistical Learning Theory, Springer, 1995.",
        "[11] Y. LeCun et al., \"Gradient-based learning applied to document recognition,\" Proc. IEEE, 1998.",
        "[13] D.C. Cirisan et al., \"Deep, big, simple neural nets for HDR,\" Neural Computation, 2010.",
        "[14] L. Wan et al., \"Regularization using DropConnect,\" Proc. ICML, 2013.",
    ]

    refs_right = [
        "[15] S. Ioffe & C. Szegedy, \"Batch normalization,\" Proc. ICML, 2015.",
        "[17] S. Sabour et al., \"Dynamic routing between capsules,\" NeurIPS, 2017.",
        "[19] K. Simonyan & A. Zisserman, \"Very deep CNNs,\" Proc. ICLR, 2015.",
        "[20] V. Nair & G.E. Hinton, \"Rectified linear units,\" Proc. ICML, 2010.",
        "[22] N. Srivastava et al., \"Dropout,\" JMLR, vol. 15, 2014.",
        "[23] D.P. Kingma & J. Ba, \"Adam optimizer,\" Proc. ICLR, 2015.",
        "[24] A. Krizhevsky et al., \"ImageNet classification with deep CNNs,\" NeurIPS, 2012.",
    ]

    ref_tf_left = add_textbox(slide10, 200000, 850000, 4400000, 3800000,
                               font_name=FONT_BODY, font_size=10)
    for ref in refs_left:
        p = ref_tf_left.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = ref
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.color.rgb = COLOR_BLACK

    ref_tf_right = add_textbox(slide10, 4700000, 850000, 4200000, 3800000,
                                font_name=FONT_BODY, font_size=10)
    for ref in refs_right:
        p = ref_tf_right.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = ref
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.color.rgb = COLOR_BLACK

    # ============================================================
    # SLIDE 11: THANK YOU
    # ============================================================
    slide11 = add_content_slide(prs, content_slide_ref)

    add_textbox(slide11, 0, 1000000, 9144000, 800000,
                text="Thank You!",
                font_name=FONT_TITLE, font_size=44, font_color=COLOR_DARK_BLUE,
                bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide11, 0, 1800000, 9144000, 400000,
                text="Questions & Discussion",
                font_name=FONT_BODY, font_size=24, font_color=COLOR_TITLE,
                bold=False, alignment=PP_ALIGN.CENTER)

    contact_tf = add_textbox(slide11, 1500000, 2600000, 6100000, 1800000,
                              font_name=FONT_BODY, font_size=14,
                              alignment=PP_ALIGN.CENTER)

    contact_lines = [
        ("Abhay Singh", True, 16, COLOR_BLACK),
        ("Department of Computer Science", False, 13, COLOR_BLACK),
        ("COER University, Roorkee, India", False, 13, COLOR_BLACK),
        ("", False, 10, COLOR_BLACK),
        ("abhaychauhan5051a@gmail.com", False, 13, COLOR_ACCENT_BLUE),
    ]

    for text, bold, size, color in contact_lines:
        p = contact_tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    add_textbox(slide11, 0, 4400000, 9144000, 500000,
                text="ICSSCS 2026 -- International Conference on Smart and Sustainable Computing Systems",
                font_name=FONT_BODY, font_size=12, font_color=COLOR_TITLE,
                bold=True, italic=True, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SAVE
    # ============================================================
    prs.save(OUTPUT_PATH)
    print(f"\n[OK] Presentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"\nSlide overview:")
    slides_info = [
        "1.  Title Slide",
        "2.  Table of Contents",
        "3.  Introduction",
        "4.  Related Work",
        "5.  Methodology (Dataset & Preprocessing)",
        "6.  Methodology (CNN Architecture & Training)",
        "7.  Result Analysis (Training & Comparison)",
        "8.  Result Analysis (Per-Class Report)",
        "9.  Conclusion & Future Work",
        "10. References",
        "11. Thank You"
    ]
    for s in slides_info:
        print(f"  {s}")


if __name__ == "__main__":
    generate_presentation()
