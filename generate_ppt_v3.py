"""
Handwritten Digit Recognition — Detailed PPT Generator v3
Sections: Title, TOC, Introduction, Dataset, Methodology (3 slides), Results (4 slides),
          Conclusion, Future Work, References, Thank You
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os, copy

BASE   = r"c:\HandWrittenDigit"
IMGS   = os.path.join(BASE, "sample_images")
TMPL   = os.path.join(BASE, "ICSSCS 2026 PPT Format.pptx")
OUT    = os.path.join(BASE, "ICSSCS_2026_Detailed_v3.pptx")

# ── palette ──────────────────────────────────────────────────────────
C_TITLE   = RGBColor(0xA5,0x30,0x10)
C_DBLUE   = RGBColor(0x1B,0x3A,0x5C)
C_ABLUE   = RGBColor(0x2E,0x75,0xB6)
C_RED     = RGBColor(0x8B,0x1A,0x1A)
C_WHITE   = RGBColor(0xFF,0xFF,0xFF)
C_BLACK   = RGBColor(0x00,0x00,0x00)
C_THEAD   = RGBColor(0x2E,0x4A,0x6E)
C_TALT    = RGBColor(0xE8,0xEF,0xF5)
C_HL      = RGBColor(0xC0,0x39,0x2B)
C_GREEN   = RGBColor(0x72,0x86,0x53)
C_DGRAY   = RGBColor(0x44,0x44,0x44)

FT = "Times New Roman"
FB = "Century Gothic"
FC = "Consolas"

# ── helpers ──────────────────────────────────────────────────────────
def tb(slide, l, t, w, h, text="", fn=FB, fs=13, fc=C_BLACK,
       bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    s = slide.shapes.add_textbox(Emu(l),Emu(t),Emu(w),Emu(h))
    tf = s.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    if text:
        r = p.add_run(); r.text = text
        r.font.name=fn; r.font.size=Pt(fs)
        r.font.color.rgb=fc; r.font.bold=bold; r.font.italic=italic
    return tf

def ap(tf, text="", fn=FB, fs=13, fc=C_BLACK, bold=False,
       italic=False, align=PP_ALIGN.LEFT, sb=0, sa=0):
    p = tf.add_paragraph(); p.alignment=align
    p.space_before=Pt(sb); p.space_after=Pt(sa)
    if text:
        r = p.add_run(); r.text=text
        r.font.name=fn; r.font.size=Pt(fs)
        r.font.color.rgb=fc; r.font.bold=bold; r.font.italic=italic
    return p

def bullets(slide, pts, l, t, w, h, fs=12, fc=C_BLACK):
    tf = tb(slide,l,t,w,h,fn=FB,fs=fs)
    first=True
    for pt in pts:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first=False; p.alignment=PP_ALIGN.LEFT
        p.space_before=Pt(2); p.space_after=Pt(2)
        rb=p.add_run(); rb.text="• "
        rb.font.name=FB; rb.font.size=Pt(fs)
        rb.font.color.rgb=C_HL; rb.font.bold=True
        rt=p.add_run(); rt.text=pt
        rt.font.name=FB; rt.font.size=Pt(fs); rt.font.color.rgb=fc
    return tf

def tbl(slide, data, l, t, w, h, col_w=None):
    nr=len(data); nc=len(data[0]) if data else 1
    shp=slide.shapes.add_table(nr,nc,Emu(l),Emu(t),Emu(w),Emu(h))
    tbl=shp.table
    if col_w:
        for i,cw in enumerate(col_w): tbl.columns[i].width=Emu(cw)
    for ri,row in enumerate(data):
        for ci,val in enumerate(row):
            cell=tbl.cell(ri,ci); cell.text=""
            p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
            r=p.add_run(); r.text=str(val)
            if ri==0:
                r.font.name=FB; r.font.size=Pt(9); r.font.bold=True
                r.font.color.rgb=C_WHITE; cell.fill.solid()
                cell.fill.fore_color.rgb=C_THEAD
            else:
                r.font.name=FB; r.font.size=Pt(8)
                r.font.color.rgb=C_BLACK; cell.fill.solid()
                cell.fill.fore_color.rgb=C_TALT if ri%2==0 else C_WHITE
            cell.margin_left=Emu(45720); cell.margin_right=Emu(45720)
            cell.margin_top=Emu(18288); cell.margin_bottom=Emu(18288)

def hdr(slide, text, l=2800000, t=80000, w=3500000, h=620000, fs=28):
    tb(slide,l,t,w,h,text,FT,fs,C_DBLUE,True,align=PP_ALIGN.CENTER)

def sec(slide, text, l, t, w=4200000, h=330000, fs=15):
    tb(slide,l,t,w,h,text,FB,fs,C_RED,True)

def line(slide, l, t, w, c=C_TITLE):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Emu(l),Emu(t),Emu(w),Pt(2))
    s.fill.solid(); s.fill.fore_color.rgb=c; s.line.fill.background()

def new_slide(prs, ref=None):
    layout = prs.slide_layouts[1]  # SECTION_HEADER_1 — matches template content slides
    return prs.slides.add_slide(layout)

def add_img(slide, fname, l, t, w):
    path=os.path.join(IMGS,fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path,Emu(l),Emu(t),Emu(w))

# ── BUILD ─────────────────────────────────────────────────────────────
def build():
    prs=Presentation(TMPL)
    ref=prs.slides[2]

    # ═══ SLIDE 1 — TITLE ═══════════════════════════════════════════
    s1=prs.slides[0]
    for sh in s1.shapes:
        if sh.has_text_frame and "Paper ID" in sh.text:
            tf=sh.text_frame
            for pa in tf.paragraphs:
                for ru in pa.runs: ru.text=""
            tf.paragraphs[0].runs[0].text="Paper ID: HDR-CNN-2026"
            tf.paragraphs[0].runs[0].font.size=Pt(14)
            tf.paragraphs[0].runs[0].font.color.rgb=C_GREEN
            p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
            r2=p2.add_run()
            r2.text="Handwritten Digit Recognition Using Convolutional Neural Networks: A Deep Learning Approach"
            r2.font.name=FB; r2.font.size=Pt(15); r2.font.color.rgb=C_GREEN; r2.font.bold=True
        if sh.has_text_frame and "Presented by" in sh.text:
            tf=sh.text_frame
            lines=[("Presented by",True,14,C_BLACK),
                   ("Abhay Singh",True,15,C_BLACK),
                   ("Dept. of Computer Science, COER University, Roorkee, India",False,13,C_BLACK),
                   ("",False,10,C_BLACK),
                   ("Date: May 2026",True,13,C_BLACK)]
            for i,(tx,bd,sz,cl) in enumerate(lines):
                if i<len(tf.paragraphs):
                    for ru in tf.paragraphs[i].runs: ru.text=""
                    if tf.paragraphs[i].runs:
                        tf.paragraphs[i].runs[0].text=tx
                        tf.paragraphs[i].runs[0].font.size=Pt(sz)
                        tf.paragraphs[i].runs[0].font.bold=bd
                        tf.paragraphs[i].runs[0].font.color.rgb=cl

    # ═══ SLIDE 2 — TABLE OF CONTENTS ═══════════════════════════════
    s2=prs.slides[1]
    toc=["1.  Introduction & Problem Statement",
         "2.  Dataset Description (MNIST)",
         "3.  Methodology — Data Preprocessing",
         "4.  Methodology — CNN Architecture",
         "5.  Methodology — Training Strategy",
         "6.  Results — Training Performance",
         "7.  Results — Test Accuracy & Per-Class Report",
         "8.  Results — Confusion Matrix & Comparison",
         "9.  Conclusion & Future Work",
         "10. References"]
    for sh in s2.shapes:
        if sh.has_text_frame and "Introduction" in sh.text:
            tf=sh.text_frame
            for i,item in enumerate(toc):
                if i<len(tf.paragraphs):
                    for ru in tf.paragraphs[i].runs: ru.text=""
                    if tf.paragraphs[i].runs:
                        tf.paragraphs[i].runs[0].text=item
                        tf.paragraphs[i].runs[0].font.size=Pt(13)
                        tf.paragraphs[i].runs[0].font.bold=True

    # ═══ SLIDE 3 — INTRODUCTION ═════════════════════════════════════
    s3=prs.slides[2]
    sec(s3,"Background & Motivation",300000,800000)
    bullets(s3,[
        "Handwritten Digit Recognition (HDR) — classify images of digits 0–9",
        "Fundamental problem in computer vision with real-world impact",
        "Postal services: auto-read ZIP codes (billions of items/year)",
        "Banking: process handwritten cheque amounts (millions/day)",
        "Healthcare: read prescription numbers & patient IDs",
        "Mobile devices: on-screen handwriting input keyboards",
    ],300000,1150000,4250000,2200000,fs=11)

    sec(s3,"Problem & Proposed Solution",4750000,800000)
    bullets(s3,[
        "Traditional ML (SVM, k-NN) needs hand-crafted features",
        "MLP flattens image → loses all spatial pixel relationships",
        "784 inputs × 512 neurons = 401K params in first layer alone!",
        "Proposed: CNN preserves 2D structure via sliding filters",
        "Automatic feature learning: edges → curves → digit shapes",
        "Deployed as real-time Streamlit web application",
    ],4750000,1150000,4150000,2200000,fs=11)

    tb(s3,200000,3500000,8700000,80000,
       "─"*120,FB,6,C_TITLE)
    tb(s3,200000,3600000,8700000,700000,
       "Abstract: A CNN trained on 60,000 MNIST images (28×28 grayscale) achieves 98.91% test accuracy. "
       "Architecture: Conv2D(32)→Conv2D(64)→MaxPool→Dropout→Dense(128)→BatchNorm→Dense(10,Softmax). "
       "Regularised with Dropout(0.25 & 0.50) and BatchNormalization. "
       "Deployed via interactive Streamlit web application with real-time canvas drawing.",
       FB,10,C_DBLUE,italic=True,align=PP_ALIGN.CENTER)

    # ═══ SLIDE 4 — DATASET ══════════════════════════════════════════
    s4=new_slide(prs,ref)
    hdr(s4,"Dataset Description — MNIST")
    line(s4,2500000,680000,4100000)

    sec(s4,"MNIST Dataset Overview",300000,800000,4200000)
    bullets(s4,[
        "Modified National Institute of Standards & Technology (1998)",
        "70,000 grayscale images of handwritten digits 0–9",
        "Training set: 60,000 images | Test set: 10,000 images",
        "Resolution: 28 × 28 pixels, single grayscale channel",
        "White digit on black background (anti-aliased edges)",
        "Digits centered by centre-of-mass, scaled to 20×20 within frame",
        "Roughly balanced: ~5,500–7,000 samples per class",
        "De facto standard benchmark — used since LeNet-5 (1998)",
    ],300000,1150000,4200000,2600000,fs=11)

    tbl(s4,[
        ["Property","Value"],
        ["Total Images","70,000"],
        ["Training Set","60,000"],
        ["Test Set","10,000"],
        ["Image Size","28 × 28 px"],
        ["Colour Space","Grayscale (1 ch)"],
        ["Pixel Range","0 (black)–255 (white)"],
        ["Classes","10 (digits 0–9)"],
        ["Format","White digit on black"],
        ["Class Balance","~5,500–7,000 / class"],
    ],4700000,800000,4100000,2400000)

    sec(s4,"Class Distribution",300000,4000000,8500000)
    tb(s4,300000,4300000,8500000,500000,
       "Digit 0: 5,923 | Digit 1: 6,742 | Digit 2: 5,958 | Digit 3: 6,131 | Digit 4: 5,842\n"
       "Digit 5: 5,421 | Digit 6: 5,918 | Digit 7: 6,265 | Digit 8: 5,851 | Digit 9: 5,949",
       FC,11,C_DBLUE,align=PP_ALIGN.CENTER)
    tb(s4,300000,4800000,8500000,300000,
       "✓ Near-uniform distribution — no significant class imbalance in MNIST",
       FB,11,C_GREEN,bold=True,align=PP_ALIGN.CENTER)

    # ═══ SLIDE 5 — METHODOLOGY: DATA PREPROCESSING ═════════════════
    s5=new_slide(prs,ref)
    hdr(s5,"Methodology — Data Preprocessing",1800000,80000,5500000)
    line(s5,1800000,680000,5500000)

    sec(s5,"Step 1 — Reshaping",200000,780000,4300000)
    tb(s5,200000,1080000,4300000,900000,
       "Raw shape:   (60000, 28, 28)   [integer pixels]\n"
       "Target shape: (60000, 28, 28, 1) [add channel dim]\n\n"
       "CNNs require 4D input: (batch, H, W, channels)\n"
       "Grayscale → channels = 1  (RGB would be 3)",
       FC,10,C_DBLUE)

    sec(s5,"Step 2 — Normalisation",200000,2000000,4300000)
    tb(s5,200000,2300000,4300000,700000,
       "x = x.astype('float32') / 255.0\n\n"
       "Maps [0, 255] integers → [0.0, 1.0] floats\n"
       "• Prevents exploding gradients during backprop\n"
       "• Ensures uniform feature scale for optimizer\n"
       "• Matches expected input range for ReLU + BatchNorm",
       FC,10,C_DBLUE)

    sec(s5,"Step 3 — One-Hot Encoding",200000,3100000,4300000)
    tb(s5,200000,3400000,4300000,700000,
       "keras.utils.to_categorical(y, num_classes=10)\n\n"
       "Label 3 → [0, 0, 0, 1, 0, 0, 0, 0, 0, 0]\n"
       "Label 7 → [0, 0, 0, 0, 0, 0, 0, 1, 0, 0]\n\n"
       "Required by Categorical Cross-Entropy loss:\n"
       "L = -Σ y_true × log(y_pred)",
       FC,10,C_DBLUE)

    sec(s5,"Pipeline Summary",4700000,780000,4200000)
    steps=[
        ("Raw MNIST","(N,28,28) uint8 [0-255]","C_HL"),
        ("Reshape","(N,28,28,1) — add channel",""),
        ("Normalize","÷ 255.0 → float32 [0,1]",""),
        ("One-Hot","Labels → 10-dim vectors",""),
        ("CNN Input","Ready for training",""),
    ]
    ty=1080000
    for i,(name,desc,_) in enumerate(steps):
        tb(s5,4700000,ty,1600000,280000,f"▶  {name}",FB,11,C_RED,bold=True)
        tb(s5,6400000,ty,2500000,280000,desc,FB,10,C_DBLUE)
        if i<4:
            tb(s5,5300000,ty+290000,400000,220000,"↓",FB,14,C_HL,bold=True,align=PP_ALIGN.CENTER)
        ty+=520000

    sec(s5,"Augmentation (Training Only)",4700000,3800000,4200000)
    bullets(s5,[
        "ImageDataGenerator applied ONLY to training batches",
        "Rotation: ±8°  |  Shift: ±8%  |  Zoom: ±8%",
        "NOT applied to validation set (avoids inflated val_loss)",
        "Balanced class weights: min(total/(10×count), 1.3×)",
    ],4700000,4100000,4200000,1100000,fs=10)

    # ═══ SLIDE 6 — METHODOLOGY: CNN ARCHITECTURE ════════════════════
    s6=new_slide(prs,ref)
    hdr(s6,"Methodology — CNN Architecture (v4-final)",1600000,80000,5900000)
    line(s6,1600000,680000,5900000)

    # Architecture diagram (text-based)
    atf=tb(s6,150000,780000,3700000,4500000,"CNN Architecture Diagram",
           FB,12,C_RED,bold=True)
    arch=[
        ("INPUT","28 × 28 × 1  grayscale image",""),
        ("","──── BLOCK 1 ────","block"),
        ("Conv2D(32, 3×3, same)","→ (28,28,32)  320 params","conv"),
        ("BatchNorm + ReLU","→ (28,28,32)  128 params","bn"),
        ("Conv2D(32, 3×3, same)","→ (28,28,32)  9,248 params","conv"),
        ("BatchNorm + ReLU","→ (28,28,32)  128 params","bn"),
        ("MaxPooling(2×2)","→ (14,14,32)  0 params","pool"),
        ("Dropout(0.20)","→ (14,14,32)  regularise","drop"),
        ("","──── BLOCK 2 ────","block"),
        ("Conv2D(64, 3×3, same)","→ (14,14,64)  18,496 params","conv"),
        ("BatchNorm + ReLU","→ (14,14,64)  256 params","bn"),
        ("Conv2D(64, 3×3, same)","→ (14,14,64)  36,928 params","conv"),
        ("BatchNorm + ReLU","→ (14,14,64)  256 params","bn"),
        ("MaxPooling(2×2)","→ (7,7,64)    0 params","pool"),
        ("Dropout(0.25)","→ (7,7,64)    regularise","drop"),
        ("","──── HEAD ────","block"),
        ("GlobalAvgPooling2D","→ (64,)        0 params","pool"),
        ("Dense(128, ReLU)","→ (128,)       8,320 params","dense"),
        ("BatchNorm+Dropout(0.35)","→ (128,)       512 params","bn"),
        ("Dense(10, Softmax)","→ (10,)        1,290 params","out"),
    ]
    color_map={"conv":C_ABLUE,"bn":C_GREEN,"pool":C_HL,
               "drop":RGBColor(0x80,0x40,0x00),"dense":C_DBLUE,
               "out":C_RED,"block":C_DGRAY,"":C_BLACK}
    for lname,lout,ltype in arch:
        p=atf.add_paragraph(); p.space_before=Pt(1); p.space_after=Pt(1)
        p.alignment=PP_ALIGN.LEFT
        r1=p.add_run(); r1.text=f"  {lname}"
        r1.font.name=FC; r1.font.size=Pt(8.5)
        r1.font.color.rgb=color_map.get(ltype,C_BLACK)
        r1.font.bold=(ltype=="block" or ltype=="out")
        if lout:
            r2=p.add_run(); r2.text=f"  {lout}"
            r2.font.name=FC; r2.font.size=Pt(8.5)
            r2.font.color.rgb=C_DGRAY

    ptf=atf.add_paragraph(); ptf.space_before=Pt(4)
    pr=ptf.add_run()
    pr.text="  ★ Total trainable params: ~76,754 (~77K)"
    pr.font.name=FB; pr.font.size=Pt(9)
    pr.font.color.rgb=C_HL; pr.font.bold=True

    # Layer explanations right side
    sec(s6,"Layer-by-Layer Purpose",4000000,780000,4900000)
    exps=[
        ("Conv2D Layers",
         "Slide 3×3 learnable filters across image. "
         "32 filters detect edges/corners; 64 filters detect curves/loops. "
         "'same' padding preserves spatial dimensions."),
        ("BatchNorm + ReLU",
         "BatchNorm: normalise activations to μ=0, σ=1 per mini-batch. "
         "Allows higher learning rates, reduces internal covariate shift. "
         "ReLU: f(x)=max(0,x) — avoids vanishing gradients."),
        ("MaxPooling2D(2×2)",
         "Retains max value in each 2×2 window. "
         "Halves spatial dims (translation invariance). "
         "Reduces computation by 75% per block."),
        ("Dropout",
         "Block 1: 20% neurons randomly zeroed. "
         "Block 2: 25% neurons randomly zeroed. "
         "Head: 35%. Forces redundant feature learning."),
        ("GlobalAvgPooling",
         "Replaces Flatten→Dense bottleneck. "
         "Computes spatial mean per feature map (7×7→1). "
         "Outputs compact 64-dim vector; fewer parameters."),
        ("Dense(10, Softmax)",
         "Output layer. Softmax: e^zi / Σe^zj. "
         "Produces probability distribution over 10 classes. "
         "Predicted digit = argmax of output vector."),
    ]
    ety=1150000
    for ename,edesc in exps:
        tb(s6,4000000,ety,4900000,250000,f"▸ {ename}",FB,11,C_RED,bold=True)
        tb(s6,4000000,ety+240000,4900000,400000,edesc,FB,9,C_DBLUE)
        ety+=660000

    # ═══ SLIDE 7 — METHODOLOGY: TRAINING STRATEGY ══════════════════
    s7=new_slide(prs,ref)
    hdr(s7,"Methodology — Training Strategy",2000000,80000,5100000)
    line(s7,2000000,680000,5100000)

    sec(s7,"Compilation Configuration",200000,780000,4300000)
    tbl(s7,[
        ["Setting","Value","Justification"],
        ["Loss Function","CategoricalCrossentropy","Standard multi-class; label_smoothing=0.05"],
        ["Optimizer","Adam (lr=1e-3)","Adaptive per-param LR; fast convergence"],
        ["Metrics","Accuracy","Standard classification metric"],
        ["Batch Size","256","Larger batches → fewer steps/epoch"],
        ["Max Epochs","15","With EarlyStopping; actual: ~9-12"],
        ["Val Split","10% (6,000 imgs)","No augmentation on validation"],
    ],200000,1100000,4300000,1800000)

    sec(s7,"Callbacks (Training Controls)",200000,3050000,4300000)
    bullets(s7,[
        "EarlyStopping: monitor val_accuracy, patience=5, restore_best_weights=True",
        "ModelCheckpoint: save best model to disk (monitor val_accuracy)",
        "ReduceLROnPlateau: halve LR when val_loss stagnates (patience=3, min_lr=1e-6)",
    ],200000,3380000,4300000,1000000,fs=11)

    sec(s7,"Augmentation Strategy (ImageDataGenerator)",200000,4480000,4300000)
    bullets(s7,[
        "rotation_range=8°  |  width/height_shift=8%  |  zoom_range=8%",
        "Applied ONLY to training flow — NOT inside model layers",
        "Reason: in-model augmentation applies to val set too (TF bug)",
        "Balanced class weights capped at 1.3× to avoid over-biasing",
    ],200000,4780000,4300000,1000000,fs=10)

    sec(s7,"Why Each Choice Matters",4650000,780000,4200000)
    choices=[
        ("Label Smoothing 0.05",
         "Prevents overconfident softmax predictions. "
         "Model learns calibrated probabilities rather than hard 0/1 targets."),
        ("Adam Optimizer",
         "Combines momentum + RMSProp. Per-parameter adaptive LR. "
         "No need for manual LR tuning. Works well with BatchNorm."),
        ("Balanced Class Weights",
         "Digit 5 has fewest samples (5,421 vs 6,742 for digit 1). "
         "Weights scale loss contribution to equalise learning across classes."),
        ("EarlyStopping + Restore",
         "Stops training when val_accuracy plateaus. "
         "restore_best_weights ensures final model = best checkpoint, not last epoch."),
        ("ReduceLROnPlateau",
         "When stuck, halving LR lets optimizer take finer steps "
         "to escape flat regions in loss landscape."),
        ("Global Avg Pooling",
         "Replaces large Flatten→Dense. Reduces params from ~590K to ~77K. "
         "Acts as structural regularizer — each feature map averages globally."),
    ]
    cy=1100000
    for cname,cdesc in choices:
        tb(s7,4650000,cy,4200000,240000,f"◆ {cname}",FB,11,C_RED,bold=True)
        tb(s7,4650000,cy+240000,4200000,350000,cdesc,FB,9,C_DBLUE)
        cy+=610000

    print("[Part 1] Slides 1-7 done — building slides 8-14...")
    return prs, ref

if __name__ == "__main__":
    build()
